# backend/routers/notebooks/export.py
import asyncio
from fastapi import APIRouter, Depends, HTTPException, Response
from sqlalchemy.orm import Session
import json
import os
import io
import zipfile
import re
from pathlib import Path
from werkzeug.utils import secure_filename
from backend.db import get_db
from backend.db.models.notebook import Notebook as DBNotebook
from backend.models import UserAuthDetails, TaskInfo
from backend.session import get_current_active_user, get_user_notebook_assets_path, build_lollms_client_from_params

from backend.settings import settings
from concurrent.futures import ThreadPoolExecutor

# Create a thread pool for blocking TTS operations
tts_executor = ThreadPoolExecutor(max_workers=4)

def _clean_text_for_tts(text: str) -> str:
    if not text: return ""
    # Remove markdown bold/italic (*) and headers (#)
    text = re.sub(r'[*#]', '', text)
    # Remove unicode emojis (Supplementary Multilingual Plane)
    text = re.sub(r'[\U00010000-\U0010ffff]', '', text)
    return text.strip()

def _extract_notebook_text(notebook) -> str:
    """Extract all readable text content from notebook tabs for TTS."""
    text_parts = []

    if notebook.title:
        text_parts.append(f"Title: {notebook.title}")
        text_parts.append("")

    for tab in notebook.tabs:
        tab_title = tab.get("title", "")
        tab_content = tab.get("content", "")
        tab_type = tab.get("type", "markdown")

        if tab_title:
            text_parts.append(tab_title)
            text_parts.append("")

        if not tab_content:
            continue

        if tab_type == "markdown":
            text_parts.append(tab_content)
            text_parts.append("")
        elif tab_type == "slides":
            try:
                import json
                slide_data = json.loads(tab_content)
                slides = slide_data.get("slides_data", [])
                for slide in slides:
                    if slide.get("title"):
                        text_parts.append(slide["title"])
                    for bullet in slide.get("bullets", []):
                        text_parts.append(bullet)
                    if slide.get("notes"):
                        text_parts.append(f"Notes: {slide['notes']}")
                    text_parts.append("")
            except Exception:
                pass
        elif tab_type == "html":
            # Strip HTML tags for plain text
            import html
            clean = re.sub(r'<[^>]+>', '', tab_content)
            clean = html.unescape(clean)
            text_parts.append(clean)
            text_parts.append("")
        elif tab_type in ("code", "youtube_script", "book_plan"):
            # Include code/script content as-is
            text_parts.append(tab_content)
            text_parts.append("")

    return "\n".join(text_parts).strip()

# Try imports for PDF generation and auto-install if missing
try:
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    import textwrap
except ImportError:
    try:
        import pipmaster as pm
        pm.install("reportlab")
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        import textwrap
    except:
        canvas = None
from ascii_colors import ASCIIColors, trace_exception

from backend.routers.files import (
    md2_to_html,              # markdown2 renderer with extras
    html_to_docx_bytes,       # HTML → python-docx mapper
    md_to_pdf_bytes,          # Markdown → PDF (images/tables/code/TOC)
    md_to_pptx_bytes,         # Markdown → PPTX (slides via ---)
    html_wrapper              # Wrap HTML body in a shell
)

router = APIRouter()

def _get_image_path(assets_path: Path, slide: dict) -> str:
    """Helper to resolve local image path from slide data"""
    if not slide.get('images'):
        return None
    
    try:
        # Get selected image or default to first
        idx = slide.get('selected_image_index', 0)
        if idx >= len(slide['images']):
            idx = 0
            
        img_data = slide['images'][idx]
        img_url = img_data.get('path', '')

        # Extract filename from /api/notebooks/{id}/assets/{filename}
        if '/assets/' in img_url:
            filename = img_url.split('/assets/')[-1]
        else:
            filename = os.path.basename(img_url)
            
        local_path = assets_path / filename

        if local_path.exists():
            return str(local_path)
    except Exception:
        pass
    return None

@router.post("/{notebook_id}/generate_video", response_model=TaskInfo)
def generate_video_endpoint(
    notebook_id: str,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Starts a background task to generate a video presentation from the slides."""
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found.")
    
    from backend.task_manager import task_manager
    from backend.tasks.notebook_tasks.slides_making import generate_presentation_video_task
    
    return task_manager.submit_task(
        name=f"Generate Video: {notebook.title}",
        target=generate_presentation_video_task,
        args=(current_user.username, notebook_id),
        description="Rendering video with TTS...",
        owner_username=current_user.username
    )

@router.get("/{notebook_id}/export")
async def export_notebook(
    notebook_id: str,
    format: str = "json",
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = (
        db.query(DBNotebook)
        .filter(
            DBNotebook.id == notebook_id,
            DBNotebook.owner_user_id == current_user.id
        )
        .first()
    )

    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found")

    export_format = format.lower()
    setting_key_format = "markdown" if export_format == "md" else export_format
    setting_key = f"export_to_{setting_key_format}_enabled"

    if not settings.get(setting_key, False):
        raise HTTPException(
            status_code=403,
            detail=f"Export to '{export_format}' is disabled by the administrator."
        )

    assets_path = get_user_notebook_assets_path(
        current_user.username,
        notebook_id
    )

    filename = f"{secure_filename(notebook.title)}.{export_format}"
    media_type = "application/octet-stream"
    file_content = b""

    try:
        # ---------------- JSON ----------------
        if export_format == "json":
            media_type = "application/json"
            file_content = json.dumps(
                {
                    "title": notebook.title,
                    "type": notebook.type,
                    "language": notebook.language,
                    "content": notebook.content,
                    "tabs": notebook.tabs,
                    "artefacts": notebook.artefacts,
                },
                indent=2
            ).encode("utf-8")

        # ---------------- DOCX ----------------
        elif export_format == "docx":
            html_parts = []

            for tab in notebook.tabs:
                if tab["type"] == "markdown" and tab.get("content"):
                    title = tab.get("title")
                    if title:
                        html_parts.append(f"<h1>{title}</h1>")
                    html_parts.append(md2_to_html(tab["content"]))

            if not html_parts:
                html_parts.append("<p>No content to export.</p>")

            html = html_wrapper(
                "\n".join(html_parts),
                title=notebook.title
            )

            file_content = html_to_docx_bytes(html)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        # ---------------- PDF (FIXED) ----------------
        elif export_format == "pdf":
            md_parts = []
            slide_images = []

            for tab in notebook.tabs:
                # Markdown tabs
                if tab["type"] == "markdown" and tab.get("content"):
                    title = tab.get("title")
                    if title:
                        md_parts.append(f"# {title}\n")
                    md_parts.append(tab["content"])
                    md_parts.append("\n\n")

                # Slide tabs
                elif tab["type"] == "slides" and tab.get("content"):
                    slides_data = json.loads(tab["content"]).get("slides_data", [])
                    for s in slides_data:
                        if s.get("title"):
                            md_parts.append(f"## {s['title']}\n")
                        for b in s.get("bullets", []):
                            md_parts.append(f"- {b}")
                        md_parts.append("\n")

                        img_path = _get_image_path(assets_path, s)
                        if img_path:
                            slide_images.append(img_path)

            content_md = "\n".join(md_parts).strip() or "# Empty notebook"

            file_content = md_to_pdf_bytes(
                content_md,
                extra_images=slide_images
            )
            media_type = "application/pdf"

        # ---------------- MARKDOWN ----------------
        elif export_format in ("md", "markdown"):
            md_parts = []

            for tab in notebook.tabs:
                # Regular markdown tabs
                if tab["type"] == "markdown" and tab.get("content"):
                    title = tab.get("title")
                    if title:
                        md_parts.append(f"# {title}\n")
                    md_parts.append(tab["content"])
                    md_parts.append("\n\n")

                # Slides converted to markdown
                elif tab["type"] == "slides" and tab.get("content"):
                    slides_data = json.loads(tab["content"]).get("slides_data", [])
                    for s in slides_data:
                        # Slide title as a secondary heading
                        if s.get("title"):
                            md_parts.append(f"## {s['title']}\n")
                        # Bullets
                        for b in s.get("bullets", []):
                            md_parts.append(f"- {b}")
                        md_parts.append("\n")
                        # Image reference if available
                        img_path = _get_image_path(assets_path, s)
                        if img_path:
                            img_filename = os.path.basename(img_path)
                            md_parts.append(f"![{s.get('title','image')}]({img_filename})\n")
                        md_parts.append("\n")

            content_md = "\n".join(md_parts).strip() or "# Empty notebook"
            file_content = content_md.encode("utf-8")
            media_type = "text/markdown"

        # ---------------- ZIP (slide images) ----------------
        elif export_format == "zip":
            bio = io.BytesIO()
            with zipfile.ZipFile(bio, "w", zipfile.ZIP_DEFLATED) as zf:
                for tab in notebook.tabs:
                    if tab["type"] == "slides" and tab.get("content"):
                        slides_data = json.loads(tab["content"]).get("slides_data", [])
                        for i, s in enumerate(slides_data):
                            img_path = _get_image_path(assets_path, s)
                            if img_path:
                                safe_title = secure_filename(
                                    s.get("title", "Untitled")
                                )[:30]
                                ext = os.path.splitext(img_path)[1] or ".png"
                                zf.write(
                                    img_path,
                                    f"Slide_{i+1:02d}_{safe_title}{ext}"
                                )
            bio.seek(0)
            file_content = bio.getvalue()
            media_type = "application/zip"

        # ---------------- PPTX ----------------
        elif export_format == "pptx":
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            prs.slide_width = Inches(13.3333)
            prs.slide_height = Inches(7.5)

            for tab in notebook.tabs:
                if tab["type"] == "slides" and tab.get("content"):
                    slides_data = json.loads(tab["content"]).get("slides_data", [])
                    for s in slides_data:
                        img_path = _get_image_path(assets_path, s)
                        layout = s.get("layout", "TitleImageBody")

                        if layout == "ImageOnly" and img_path:
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            slide.shapes.add_picture(
                                img_path,
                                0,
                                0,
                                width=prs.slide_width,
                                height=prs.slide_height
                            )
                        else:
                            slide = prs.slides.add_slide(prs.slide_layouts[1])
                            slide.shapes.title.text = s.get("title", "")

                        if s.get("notes"):
                            slide.notes_slide.notes_text_frame.text = s["notes"]

            bio = io.BytesIO()
            prs.save(bio)
            file_content = bio.getvalue()
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        # ---------------- AUDIO (TTS) ----------------
        elif export_format in ("wav", "audio"):
            loop = asyncio.get_running_loop()
            try:
                lc = await loop.run_in_executor(
                    tts_executor,
                    lambda: build_lollms_client_from_params(
                        username=current_user.username,
                        load_llm=False,
                        load_tts=True
                    )
                )

                if not lc.tts:
                    raise HTTPException(
                        status_code=400,
                        detail="Text-to-Speech (TTS) is not configured for this user."
                    )

                # Extract all text from the notebook
                full_text = _extract_notebook_text(notebook)

                if not full_text:
                    raise HTTPException(
                        status_code=400,
                        detail="No text content found in notebook to convert to speech."
                    )

                # Clean text for TTS
                cleaned_text = _clean_text_for_tts(full_text)

                # Determine voice and language
                voice_to_use = None
                language_to_use = None

                # Check for user's active custom voice
                db_user = db.query(DBUser).filter(DBUser.id == current_user.id).first()
                if db_user and db_user.active_voice_id:
                    from backend.db.models.voice import UserVoice as DBUserVoice
                    from backend.session import get_user_data_root
                    active_voice = db.query(DBUserVoice).filter(DBUserVoice.id == db_user.active_voice_id).first()
                    if active_voice:
                        user_voices_path = get_user_data_root(current_user.username) / "voices"
                        voice_file_path = user_voices_path / Path(active_voice.file_path)
                        if voice_file_path.exists():
                            voice_to_use = str(voice_file_path.resolve())
                            language_to_use = active_voice.language

                # Fallback to user's language preference
                if not language_to_use and db_user and db_user.ai_response_language and db_user.ai_response_language.lower() != "auto":
                    language_to_use = db_user.ai_response_language

                # Final fallback to English
                if not language_to_use:
                    language_to_use = 'en'

                # Determine model to use
                model_to_use = None
                user_tts_model_full = current_user.tts_binding_model_name
                if user_tts_model_full and '/' in user_tts_model_full:
                    _, model_name = user_tts_model_full.split('/', 1)
                    model_to_use = model_name

                def _generate_audio():
                    return lc.tts.generate_audio(
                        text=cleaned_text,
                        voice=voice_to_use,
                        model=model_to_use,
                        language=language_to_use
                    )

                audio_bytes = await loop.run_in_executor(tts_executor, _generate_audio)

                file_content = audio_bytes
                media_type = "audio/wav"
                filename = f"{secure_filename(notebook.title)}.wav"

            except HTTPException as e:
                raise e
            except Exception as e:
                trace_exception(e)
                raise HTTPException(
                    status_code=500,
                    detail=f"TTS generation failed: {str(e)}"
                )

        else:
            raise HTTPException(status_code=400, detail="Unsupported export format")

    except ImportError as e:
        raise HTTPException(
            status_code=501,
            detail=f"Missing library for '{export_format}': {e.name}"
        )
    except Exception as e:
        trace_exception(e)
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate export: {e}"
        )

    headers = {
        "Content-Disposition": f'attachment; filename="{filename}"'
    }

    return Response(
        content=file_content,
        media_type=media_type,
        headers=headers
    )
