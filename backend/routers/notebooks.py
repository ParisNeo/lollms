# backend/routers/notebooks.py
import uuid
import json
import shutil
import re
import base64
import os
from datetime import datetime
from typing import List, Optional, Dict, Any
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Body, Response
from fastapi.responses import FileResponse 
from sqlalchemy.orm import Session
from pydantic import BaseModel
from werkzeug.utils import secure_filename

from backend.db import get_db
from backend.db.models.notebook import Notebook as DBNotebook
from backend.models import UserAuthDetails, TaskInfo
from backend.session import get_current_active_user, get_user_notebook_assets_path, get_user_lollms_client
from backend.task_manager import task_manager
from ascii_colors import trace_exception
from backend.tasks.notebook_tasks import _process_notebook_task, _ingest_notebook_sources_task, _convert_file_with_docling_task

router = APIRouter(
    prefix="/api/notebooks",
    tags=["Notebooks"],
    dependencies=[Depends(get_current_active_user)]
)

class StructureItem(BaseModel):
    title: str
    type: str = "markdown" 
    content: Optional[str] = ""

class NotebookCreate(BaseModel):
    title: str
    content: Optional[str] = ""
    type: Optional[str] = "generic"
    structure: Optional[List[StructureItem]] = None 
    initialPrompt: Optional[str] = None
    urls: Optional[List[str]] = None
    youtube_urls: Optional[List[str]] = None
    raw_text: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None 

class NotebookUpdate(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None
    artefacts: Optional[List[Dict[str, Any]]] = None
    tabs: Optional[List[Dict[str, Any]]] = None

class NotebookResponse(BaseModel):
    id: str
    title: str
    content: str
    type: str
    artefacts: List[Dict[str, Any]]
    tabs: List[Dict[str, Any]]
    created_at: datetime
    updated_at: datetime
    class Config: from_attributes = True

class GenerateStructureRequest(BaseModel):
    type: str
    prompt: str
    urls: Optional[List[str]] = []
    files: Optional[List[Any]] = [] 

class GenerateTitleResponse(BaseModel):
    title: str

class ProcessRequest(BaseModel):
    prompt: str
    input_tab_ids: List[str]
    output_type: str
    target_tab_id: Optional[str] = None
    skip_llm: bool = False
    generate_speech: bool = False

@router.get("", response_model=List[NotebookResponse])
def get_notebooks(
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    return db.query(DBNotebook).filter(DBNotebook.owner_user_id == current_user.id).order_by(DBNotebook.updated_at.desc()).all()

@router.post("", response_model=NotebookResponse)
def create_notebook(
    payload: NotebookCreate,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    initial_tabs = []
    if payload.structure:
        for item in payload.structure:
            initial_tabs.append({
                "id": str(uuid.uuid4()),
                "title": item.title,
                "type": item.type,
                "content": item.content,
                "images": []
            })
    
    if not initial_tabs:
         initial_tabs.append({
            "id": str(uuid.uuid4()),
            "title": "Main",
            "type": "markdown",
            "content": payload.initialPrompt or "",
            "images": []
        })

    content_to_store = payload.content or ""
    if payload.metadata:
        try:
            content_obj = { "text": payload.content or "", "metadata": payload.metadata }
            content_to_store = json.dumps(content_obj)
        except: pass

    new_notebook = DBNotebook(
        title=payload.title,
        content=content_to_store,
        type=payload.type,
        owner_user_id=current_user.id,
        tabs=initial_tabs,
        artefacts=[]
    )
    
    if payload.raw_text:
        new_notebook.artefacts = [{
            "filename": "Initial Notes",
            "content": payload.raw_text,
            "type": "text",
            "is_loaded": True
        }]

    db.add(new_notebook)
    db.commit()
    db.refresh(new_notebook)
    
    if (payload.urls and len(payload.urls) > 0) or (payload.youtube_urls and len(payload.youtube_urls) > 0):
        task_manager.submit_task(
            name=f"Notebook Ingest: {new_notebook.title}",
            target=_ingest_notebook_sources_task,
            args=(current_user.username, new_notebook.id, payload.urls or [], payload.youtube_urls or []),
            description="Scraping web and YouTube sources...",
            owner_username=current_user.username
        )

    return new_notebook

@router.put("/{notebook_id}", response_model=NotebookResponse)
def update_notebook(
    notebook_id: str,
    payload: NotebookUpdate,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found.")
    
    if payload.title is not None: notebook.title = payload.title
    if payload.content is not None: notebook.content = payload.content
    if payload.artefacts is not None: notebook.artefacts = payload.artefacts
    if payload.tabs is not None: notebook.tabs = payload.tabs
    
    notebook.updated_at = datetime.now()
    db.commit()
    db.refresh(notebook)
    return notebook

@router.delete("/{notebook_id}")
def delete_notebook(
    notebook_id: str,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found.")
    
    assets_path = get_user_notebook_assets_path(current_user.username, notebook_id)
    if assets_path.exists():
        shutil.rmtree(assets_path, ignore_errors=True)
        
    db.delete(notebook)
    db.commit()
    return {"message": "Notebook deleted."}

@router.post("/structure", response_model=List[StructureItem])
def generate_notebook_structure(
    request: GenerateStructureRequest,
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    lc = get_user_lollms_client(current_user.username)
    prompt = f"""Create a structure for a '{request.type}' notebook based on: "{request.prompt}".
    Return a JSON list of objects with "title", "type" (markdown/slides), and "content".
    """
    try:
        response_text = lc.generate_text(prompt, max_new_tokens=1024)
        match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if match:
            return [StructureItem(**item) for item in json.loads(match.group(0))]
        return [StructureItem(title="Chapter 1", content=response_text)]
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/{notebook_id}/upload")
async def upload_notebook_source(
    notebook_id: str,
    file: UploadFile = File(...),
    use_docling: bool = Form(False),
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found.")

    assets_path = get_user_notebook_assets_path(current_user.username, notebook_id)
    assets_path.mkdir(parents=True, exist_ok=True)
    
    fn = secure_filename(file.filename)
    unique_fn = f"{uuid.uuid4().hex[:8]}_{fn}"
    file_path = assets_path / unique_fn
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    content = ""
    is_text = file_path.suffix.lower() in ['.txt', '.md', '.py', '.js', '.json', '.html', '.css', '.c', '.cpp']
    if is_text:
        content = file_path.read_text(encoding='utf-8', errors='ignore')
        new_art = { "filename": unique_fn, "content": content, "type": "text", "is_loaded": True }
        notebook.artefacts = list(notebook.artefacts) + [new_art]
        db.commit()
    elif use_docling:
        task_manager.submit_task(
            name=f"Convert: {fn}",
            target=_convert_file_with_docling_task,
            args=(current_user.username, notebook_id, str(file_path), unique_fn),
            owner_username=current_user.username
        )
    
    return {"filename": unique_fn}

@router.post("/{notebook_id}/describe_image")
async def describe_notebook_image(
    notebook_id: str,
    file: UploadFile = File(...),
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    try:
        content = await file.read()
        b64 = base64.b64encode(content).decode('utf-8')
        lc = get_user_lollms_client(current_user.username)
        desc = lc.generate_text("Describe this image for a prompt:", images=[b64])
        return {"description": desc}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/{notebook_id}/describe_asset")
async def describe_notebook_asset_endpoint(
    notebook_id: str,
    payload: Dict[str, str],
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    fn = payload.get("filename")
    path = get_user_notebook_assets_path(current_user.username, notebook_id) / secure_filename(fn)
    if not path.exists():
        raise HTTPException(status_code=404)
    b64 = base64.b64encode(path.read_bytes()).decode('utf-8')
    lc = get_user_lollms_client(current_user.username)
    return {"description": lc.generate_text("Describe this image for a prompt:", images=[b64])}

@router.post("/{notebook_id}/artefact")
def create_text_artefact(
    notebook_id: str,
    payload: Dict[str, str], 
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook: raise HTTPException(status_code=404)
    fn = secure_filename(payload.get('title', 'note.txt'))
    path = get_user_notebook_assets_path(current_user.username, notebook_id) / fn
    path.write_text(payload.get('content', ''), encoding='utf-8')
    notebook.artefacts = list(notebook.artefacts) + [{ "filename": fn, "content": payload.get('content', ''), "type": "text", "is_loaded": True }]
    db.commit()
    return {"filename": fn}

@router.post("/{notebook_id}/generate_title", response_model=GenerateTitleResponse)
def generate_notebook_title(
    notebook_id: str,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook: raise HTTPException(status_code=404)
    lc = get_user_lollms_client(current_user.username)
    title = lc.generate_text(f"Summarize this notebook content into a short title: {notebook.content[:1000]}").strip().strip('"')
    notebook.title = title
    db.commit()
    return {"title": title}

@router.post("/{notebook_id}/scrape")
def scrape_url_to_notebook(
    notebook_id: str,
    payload: Dict[str, str],
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    return task_manager.submit_task(
        name=f"Scrape: {payload['url']}",
        target=_ingest_notebook_sources_task,
        args=(current_user.username, notebook_id, [payload['url']], []),
        owner_username=current_user.username
    )

@router.post("/{notebook_id}/process", response_model=TaskInfo)
def process_notebook_ai(
    notebook_id: str,
    payload: ProcessRequest,
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    return task_manager.submit_task(
        name=f"AI Task: {payload.output_type}",
        target=_process_notebook_task,
        args=(current_user.username, notebook_id, payload.prompt, payload.input_tab_ids, payload.output_type, payload.target_tab_id, payload.skip_llm, payload.generate_speech),
        owner_username=current_user.username
    )

@router.get("/{notebook_id}/export")
def export_notebook(
    notebook_id: str,
    format: str = "json",
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    notebook = db.query(DBNotebook).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook not found.")

    if format == "json":
        data = { "title": notebook.title, "type": notebook.type, "content": notebook.content, "tabs": notebook.tabs, "artefacts": notebook.artefacts }
        return Response(content=json.dumps(data, indent=2), media_type="application/json", headers={"Content-Disposition": f"attachment; filename={secure_filename(notebook.title)}.json"})
    
    elif format == "pdf":
        try:
            from markdown_pdf import MarkdownPdf, Section
            import tempfile
            
            pdf = MarkdownPdf(toc_level=2)
            content_md = f"# {notebook.title}\n\n"
            for tab in notebook.tabs:
                content_md += f"## {tab['title']}\n\n"
                if tab['type'] == 'slides':
                    try:
                        data = json.loads(tab['content'])
                        for s in data.get('slides_data', []):
                            content_md += f"### {s['title']}\n"
                            for b in s.get('bullets', []): content_md += f"* {b}\n"
                            content_md += "\n"
                    except: pass
                else: content_md += (tab.get('content') or "") + "\n\n"
            
            pdf.add_section(Section(content_md))
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
                out_path = tf.name
            pdf.save(out_path)
            with open(out_path, "rb") as f: data = f.read()
            os.unlink(out_path)
            return Response(content=data, media_type="application/pdf", headers={"Content-Disposition": f"attachment; filename={secure_filename(notebook.title)}.pdf"})
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    elif format == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import io
            
            prs = Presentation()
            prs.slide_width = Inches(13.3333)
            prs.slide_height = Inches(7.5)
            assets_path = get_user_notebook_assets_path(current_user.username, notebook_id)
            
            for tab in notebook.tabs:
                if tab['type'] == 'slides' and tab.get('content'):
                    try:
                        slides_data = json.loads(tab['content']).get('slides_data', [])
                        for s in slides_data:
                            # Resolve Image
                            img_path = None
                            if s.get('images'):
                                img_url = s['images'][s.get('selected_image_index', 0)]['path']
                                local_fn = os.path.basename(img_url)
                                if (assets_path / local_fn).exists(): img_path = str(assets_path / local_fn)

                            # Layout mapping
                            layout = s.get('layout', 'TitleBody')
                            if layout == 'ImageOnly' and img_path:
                                slide = prs.slides.add_slide(prs.slide_layouts[6])
                                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                            else:
                                slide = prs.slides.add_slide(prs.slide_layouts[1])
                                if slide.shapes.title: slide.shapes.title.text = s.get('title', '')
                                if len(slide.placeholders) > 1:
                                    slide.placeholders[1].text = "\n".join(s.get('bullets', []))
                                if img_path: slide.shapes.add_picture(img_path, Inches(8), Inches(1.5), width=Inches(4.5))
                            
                            if s.get('speech'):
                                try: slide.notes_slide.notes_text_frame.text = s['speech']
                                except: pass
                    except: pass
            
            bio = io.BytesIO()
            prs.save(bio)
            return Response(content=bio.getvalue(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename={secure_filename(notebook.title)}.pptx"})
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    raise HTTPException(status_code=400)

@router.get("/{notebook_id}/assets/{filename}")
def get_notebook_asset(
    notebook_id: str,
    filename: str,
    current_user: UserAuthDetails = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    nb = db.query(DBNotebook.id).filter(DBNotebook.id == notebook_id, DBNotebook.owner_user_id == current_user.id).first()
    if not nb: raise HTTPException(status_code=403)
    path = get_user_notebook_assets_path(current_user.username, notebook_id) / secure_filename(filename)
    if not path.exists(): raise HTTPException(status_code=404)
    return FileResponse(path)
