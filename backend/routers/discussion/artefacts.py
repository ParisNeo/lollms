# backend/routers/discussion/artefacts.py
# Standard Library Imports
import base64
import io
import asyncio
import json
import requests
import re
import os
from datetime import datetime
from pathlib import Path
from urllib.parse import urlparse, unquote
from typing import List, Optional, Dict, Any
import pipmaster as pm
try:
    from docx import Document as DocxDocument
    from pptx import Presentation
    import pandas as pd
    import fitz
    from docx2python import docx2python
except ImportError:
    pd = None
    docx2python = None
from fastapi import APIRouter, Depends, File, HTTPException, Query, UploadFile, Form, status, Body
from sqlalchemy.orm import Session
from ascii_colors import trace_exception

from backend.db import get_db
from backend.models import UserAuthDetails, ArtefactInfo, ArtefactCreateManual, ArtefactUpdate, ExportContextRequest, LoadArtefactRequest, TaskInfo, UnloadArtefactRequest, UrlImportRequest, ArtefactAndDataZoneUpdateResponse
from backend.models.discussion import ArtefactUploadResponse
from backend.session import get_current_active_user, get_user_lollms_client
from backend.discussion import get_user_discussion
from backend.routers.discussion.helpers import get_discussion_and_owner_for_request
from backend.task_manager import task_manager
from backend.tasks.artefact_tasks import _import_artefact_from_url_task
from backend.tasks.utils import _to_task_info
from backend.routers.discussion.helpers import get_discussion_and_owner_for_request
# .msg handling
try:
    import extract_msg  # pip install extract-msg
except ImportError:
    extract_msg = None

# YouTube Transcript handling
try:
    from youtube_transcript_api import YouTubeTranscriptApi
except ImportError:
    YouTubeTranscriptApi = None

from backend.db import get_db
from backend.db.models.user import User as DBUser
from pydantic import BaseModel

# --- New Models for Search/Select ---
class WikipediaSearchRequest(BaseModel):
    query: str

class WikipediaImportItem(BaseModel):
    title: str
    url: str

class WikipediaImportSelectedRequest(BaseModel):
    items: List[WikipediaImportItem]
    auto_load: bool = True

class ArxivSearchRequest(BaseModel):
    query: Optional[str] = None
    author: Optional[str] = None
    year: Optional[int] = None
    max_results: int = 5

class ArxivImportItem(BaseModel):
    id: str
    title: str
    mode: str = "abstract" # "abstract" or "full"

class ArxivImportSelectedRequest(BaseModel):
    items: List[ArxivImportItem]
    auto_load: bool = True

class WebSearchRequest(BaseModel):
    query: str
    provider: str

class WikipediaImportRequest(BaseModel):
    query: str
    auto_load: bool = True

class YoutubeTranscriptImportRequest(BaseModel):
    video_url: str
    language: str = "en"
    auto_load: bool = True

class GithubImportRequest(BaseModel):
    url: str
    auto_load: bool = True

class GithubSearchRequest(BaseModel):
    query: str

class StackOverflowImportRequest(BaseModel):
    url: str
    auto_load: bool = True

class StackOverflowSearchRequest(BaseModel):
    query: str
class StackOverflowImportRequest(BaseModel):
    video_url: str
    language: str = "en"
    auto_load: bool = True

def _map_artefact_for_ui(art: dict, discussion_id: str = None) -> dict:
    """
    Standardizes Artefact metadata for the UI.
    Maps internal library keys ('type', 'active') to public API keys ('artefact_type', 'is_loaded').
    """
    mapped = {k: v for k, v in art.items() if k not in ['content', 'images']}
    
    # Ensure discussion_id is preserved
    if 'discussion_id' not in mapped and discussion_id:
        mapped['discussion_id'] = discussion_id

    # Map library internal 'type' to UI 'artefact_type'
    mapped['artefact_type'] = art.get('type', 'document')
    
    # Map library 'active' (boolean) to UI 'is_loaded'
    mapped['is_loaded'] = bool(art.get('active', False))
    
    # Handle serialization of dates
    for date_key in ['created_at', 'updated_at']:
        if isinstance(mapped.get(date_key), datetime):
            mapped[date_key] = mapped[date_key].isoformat()
            
    return mapped

def build_artefacts_router(router: APIRouter):
    # Ensure Arxiv is available
    try:
        import arxiv
    except ImportError:
        pm.ensure_packages("arxiv")
        import arxiv

    @router.post("/{discussion_id}/artefacts/web/search", response_model=List[Dict[str, Any]])
    async def search_web(
        discussion_id: str,
        request: WebSearchRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        results =[]
        try:
            if request.provider == "duckduckgo":
                pm.ensure_packages("duckduckgo-search") # using duckduckgo-search to avoid ddgs rename warning if possible or handle it
                try:
                    from duckduckgo_search import DDGS
                except ImportError:
                    from ddgs import DDGS
                with DDGS() as ddgs:
                    raw_results =[r for r in ddgs.text(request.query, max_results=10)]
                    results =[{'title': r.get('title'), 'url': r.get('href'), 'snippet': r.get('body')} for r in raw_results]
            elif request.provider == "google":
                pm.ensure_packages("google-api-python-client")
                from googleapiclient.discovery import build as google_build
                
                db_user = db.query(DBUser).filter(DBUser.username == current_user.username).first()
                if not db_user or not db_user.google_api_key or not db_user.google_cse_id:
                    raise HTTPException(status_code=400, detail="Google Search API is not configured in settings.")
                    
                service = google_build("customsearch", "v1", developerKey=db_user.google_api_key)
                res = service.cse().list(q=request.query, cx=db_user.google_cse_id, num=10).execute()
                items = res.get('items', [])
                results =[{'title': item.get('title'), 'url': item.get('link'), 'snippet': item.get('snippet')} for item in items]
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Search failed: {e}")
        
        return results

     # safe_store is needed for RAG callbacks
    @router.get("/{discussion_id}/artefacts", response_model=List[ArtefactInfo])
    async def list_discussion_artefacts(
        discussion_id: str,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)  
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db)
        if not discussion:
            raise HTTPException(status_code=404, detail="Discussion not found")
        
        # Fetch raw artefacts from the manager and use centralized mapper
        raw_artefacts = discussion.list_artefacts()
        return [_map_artefact_for_ui(art, discussion_id) for art in raw_artefacts]

    @router.post("/{discussion_id}/artefacts", response_model=ArtefactUploadResponse)
    async def add_discussion_artefact(
        discussion_id: str,
        file: UploadFile = File(...),
        extract_images: bool = Form(True),
        auto_load: bool = Form(True),
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db) 
    ):
        print(f"extract_images: {extract_images}")
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        if not discussion:
            raise HTTPException(status_code=404, detail="Discussion not found")
        
        try:
            content_bytes = await file.read()
            title = file.filename
            extension = Path(title).suffix.lower()
            
            content = ""
            images: List[str] = []
            
            image_mimetypes = ["image/jpeg", "image/png", "image/gif", "image/webp"]
            if file.content_type in image_mimetypes:
                images.append(base64.b64encode(content_bytes).decode('utf-8'))
                content = ""
            else:
                CODE_EXTENSIONS = {
                    ".py": "python", ".js": "javascript", ".ts": "typescript", ".html": "html", ".css": "css",
                    ".c": "c", ".cpp": "cpp", ".h": "cpp", ".hpp": "cpp", ".cs": "csharp", ".java": "java",
                    ".json": "json", ".xml": "xml", ".sh": "bash", ".md": "markdown", ".vhd": "vhdl", ".v": "verilog",
                    ".rb": "ruby", ".php": "php", ".go": "go", ".rs": "rust", ".swift": "swift", ".kt": "kotlin"
                }

                if extension == ".pdf":
                    # --- NEW: High-Fidelity Multi-Modal PDF Ingestion ---
                    pdf_doc = fitz.open(stream=content_bytes, filetype="pdf")
                    pages_md = []
                    rendered_page_images = []
                    
                    for i, page in enumerate(pdf_doc):
                        # 1. Render page to image for the vision model
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("png")
                        rendered_page_images.append(base64.b64encode(img_bytes).decode('utf-8'))
                        
                        # 2. Extract text from page
                        page_text = page.get_text().strip()
                        
                        # 3. Create multimodal section with anchor
                        # The ID format matches Title::PageIndex for resolution
                        pages_md.append(
                            f"## Page {i+1}\n\n"
                            f'<artefact_image id="{title}::{i}" />\n\n'
                            f"{page_text}"
                        )
                    
                    content = "\n\n---\n\n".join(pages_md)
                    images = rendered_page_images
                    pdf_doc.close()

                elif extension == ".docx":
                    if docx2python is None:
                        content = "Error: docx2python library is not installed, so tables cannot be extracted."
                    else:
                        with io.BytesIO(content_bytes) as docx_io:
                            result = docx2python(docx_io)
                            content = result.text
                            if extract_images and result.images:
                                for image_bytes in result.images.values():
                                    images.append(base64.b64encode(image_bytes).decode("utf-8"))

                elif extension == ".xlsx" or "spreadsheetml" in file.content_type:
                    try:
                        xls = pd.read_excel(io.BytesIO(content_bytes), sheet_name=None)
                        md_parts = []
                        for sheet_name, df in xls.items():
                            md_parts.append(f"### {sheet_name}\n\n{df.to_markdown(index=False)}")
                        content = "\n\n".join(md_parts)
                    except Exception as e:
                        trace_exception(e)
                        content = f"Error processing XLSX file: {e}"

                elif extension == ".pptx":
                    try:
                        with io.BytesIO(content_bytes) as pptx_io:
                            prs = Presentation(pptx_io)
                            slide_texts: List[str] = []
                            for idx, slide in enumerate(prs.slides, start=1):
                                slide_parts: List[str] = []
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        txt = (shape.text or "").strip()
                                        if txt:
                                            slide_parts.append(txt)
                                if slide_parts:
                                    slide_texts.append(f"--- Slide {idx} ---\n" + "\n".join(slide_parts))
                            content = "\n\n".join(slide_texts)
                            if extract_images:                            
                                from pptx.enum.shapes import MSO_SHAPE_TYPE
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                            image_blob = shape.image.blob
                                            images.append(base64.b64encode(image_blob).decode("utf-8"))
                    except Exception as e:
                        trace_exception(e)
                        content = f"Error processing PPTX file: {e}"

                elif extension == ".msg":
                    if extract_msg is None:
                        content = "Error: extract-msg library is not installed."
                    else:
                        try:
                            with io.BytesIO(content_bytes) as bio:
                                import tempfile, os
                                with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tf:
                                    tf.write(bio.getvalue())
                                    temp_path = tf.name
                            try:
                                msg = extract_msg.Message(temp_path)
                                msg_sender = getattr(msg, "sender", None) or getattr(msg, "from_", None) or ""
                                msg_to = getattr(msg, "to", "") or ""
                                msg_cc = getattr(msg, "cc", "") or ""
                                msg_date = getattr(msg, "date", "")
                                msg_subject = getattr(msg, "subject", "") or title
                                msg_body = (getattr(msg, "body", "") or "").strip()
                                html_body = getattr(msg, "htmlBody", None)
                                if not msg_body and html_body:
                                    try:
                                        from bs4 import BeautifulSoup
                                        msg_body = BeautifulSoup(html_body, "html.parser").get_text()
                                    except Exception:
                                        pass

                                header_lines = []
                                if msg_subject: header_lines.append(f"# {msg_subject}")
                                meta = []
                                if msg_sender: meta.append(f"From: {msg_sender}")
                                if msg_to: meta.append(f"To: {msg_to}")
                                if msg_cc: meta.append(f"CC: {msg_cc}")
                                if msg_date: meta.append(f"Date: {msg_date}")
                                if meta: header_lines.append("\n".join(meta))
                                header = "\n\n".join(header_lines)

                                attachment_text_parts: List[str] = []
                                for att in msg.attachments:
                                    att_name = getattr(att, "longFilename", None) or getattr(att, "shortFilename", None) or "attachment"
                                    att_bytes = att.data or b""
                                    att_ext = Path(att_name).suffix.lower()

                                    if extract_images:
                                        if att_ext in [".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"]:
                                            images.append(base64.b64encode(att_bytes).decode("utf-8"))
                                            continue

                                    try:
                                        text_guess = att_bytes.decode("utf-8")
                                    except UnicodeDecodeError:
                                        try:
                                            text_guess = att_bytes.decode("latin-1", errors="replace")
                                        except Exception:
                                            text_guess = ""

                                    if att_ext in [".txt", ".md", ".py", ".json", ".csv", ".log", ".yaml", ".yml", ".xml", ".html", ".js", ".ts", ".css"]:
                                        lang = {
                                            ".md": "markdown", ".py": "python", ".json": "json", ".csv": "csv",
                                            ".yaml": "yaml", ".yml": "yaml", ".xml": "xml", ".html": "html",
                                            ".js": "javascript", ".ts": "typescript", ".css": "css"
                                        }.get(att_ext, "")
                                        fence = f"````{lang}\n{text_guess}\n````"
                                        attachment_text_parts.append(f"### Attachment: {att_name}\n\n{fence}")
                                    else:
                                        attachment_text_parts.append(f"- Attachment: {att_name} ({len(att_bytes)} bytes)")

                                attachments_md = "\n\n".join(attachment_text_parts)
                                content = "\n\n".join([p for p in [header, msg_body, attachments_md] if p])

                            finally:
                                os.unlink(temp_path)
                        except Exception as e:
                            trace_exception(e)
                            content = f"Error processing MSG file: {e}"

                elif extension in CODE_EXTENSIONS:
                    lang = CODE_EXTENSIONS[extension]
                    text_content = content_bytes.decode('utf-8', errors='replace')
                    content = f"````{lang}\n{text_content}\n````"
                else:
                    try:
                        content = content_bytes.decode('utf-8')
                    except UnicodeDecodeError:
                        content = content_bytes.decode('latin-1', errors='replace')

            # Determine type: Code stays code, everything else is a 'file' (custom registered type)
            target_type = "code" if extension in ['.py', '.js', '.ts', '.html', '.css'] else "file"

            # Check if an artefact with this title already exists in the discussion
            existing_artefact = discussion.get_artefact(title=title)

            if existing_artefact:
                # Increment version
                artefact_info = discussion.update_artefact(
                    title,
                    content,
                    new_images=images,
                    author=current_user.username,
                    active=auto_load,
                    artefact_type=target_type
                )
            else:
                # Initial creation (Version 1)
                artefact_info = discussion.add_artefact(
                    title,
                    content,
                    images=images,
                    author=current_user.username,
                    active=auto_load,
                    artefact_type=target_type
                )
            
            discussion.commit()

            all_images_info = discussion.get_discussion_images()
            
            if isinstance(artefact_info.get('created_at'), datetime):
                artefact_info['created_at'] = artefact_info['created_at'].isoformat()
            if isinstance(artefact_info.get('updated_at'), datetime):
                artefact_info['updated_at'] = artefact_info['updated_at'].isoformat()

            return {
                "new_artefact_info": artefact_info,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Failed to add artefact: {e}")

    @router.post("/{discussion_id}/artefacts/wikipedia/search", response_model=List[Dict[str, str]])
    async def search_wikipedia(
        discussion_id: str,
        request: WikipediaSearchRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        try:
            query = request.query.strip()
            lang = "en"
            
            # Check if it's a URL to just return that specific one
            if query.startswith("http://") or query.startswith("https://"):
                parsed = urlparse(query)
                if "wikipedia.org" in parsed.netloc:
                    title = unquote(parsed.path.split('/')[-1]).replace('_', ' ')
                    return [{"title": title, "url": query, "snippet": "Direct URL import."}]

            api_url = f"https://{lang}.wikipedia.org/w/api.php"
            params = {
                "action": "query",
                "list": "search",
                "srsearch": query,
                "format": "json",
                "srlimit": 10
            }
            headers = {
                "User-Agent": "LoLLMs/1.0 (https://github.com/ParisNeo/lollms; parisneo@gmail.com)"
            }
            resp = requests.get(api_url, params=params, headers=headers)
            resp.raise_for_status()
            data = resp.json()
            
            results = []
            for item in data.get("query", {}).get("search", []):
                results.append({
                    "title": item["title"],
                    "url": f"https://{lang}.wikipedia.org/wiki/{item['title'].replace(' ', '_')}",
                    "snippet": item["snippet"].replace('<span class="searchmatch">', '').replace('</span>', '')
                })
            return results
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Wikipedia search failed: {e}")

    @router.post("/{discussion_id}/artefacts/wikipedia/import", response_model=ArtefactAndDataZoneUpdateResponse)
    async def import_wikipedia_selected(
        discussion_id: str,
        request: WikipediaImportSelectedRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        try:
            for item in request.items:
                api_url = "https://en.wikipedia.org/w/api.php"
                params = {
                    "action": "query",
                    "prop": "extracts",
                    "titles": item.title,
                    "explaintext": 1,
                    "format": "json"
                }
                headers = {
                    "User-Agent": "LoLLMs/1.0 (https://github.com/ParisNeo/lollms; parisneo_ai@gmail.com)"
                }
                resp = requests.get(api_url, params=params, headers=headers)
                data = resp.json()
                pages = data.get("query", {}).get("pages", {})
                page_id = list(pages.keys())[0]
                if page_id != "-1":
                    content = pages[page_id].get("extract", "")
                    full_md = f"# {item.title}\nSource: {item.url}\n\n{content}"
                    # Use update_artefact to handle creation/versioning and activation in one go
                    discussion.update_artefact(
                        f"{item.title}.md", 
                        full_md, 
                        author=current_user.username,
                        active=request.auto_load
                    )
            
            discussion.commit()
            
            artefacts = discussion.list_artefacts()
            for artefact in artefacts:
                if isinstance(artefact.get('created_at'), datetime): artefact['created_at'] = artefact['created_at'].isoformat()
                if isinstance(artefact.get('updated_at'), datetime): artefact['updated_at'] = artefact['updated_at'].isoformat()

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            all_images_info = discussion.get_discussion_images()

            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Wikipedia import failed: {e}")

    @router.post("/{discussion_id}/artefacts/arxiv/search", response_model=List[Dict[str, Any]])
    async def search_arxiv(
        discussion_id: str,
        request: ArxivSearchRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        import arxiv
        try:
            query_parts = []
            if request.query: query_parts.append(request.query)
            if request.author: query_parts.append(f"au:{request.author}")
            
            # Construct final query string
            query_str = " AND ".join(query_parts) if query_parts else "all:*"
            
            search = arxiv.Search(
                query=query_str,
                max_results=request.max_results,
                sort_by=arxiv.SortCriterion.Relevance
            )
            
            results = []
            for r in search.results():
                # Filter by year client-side if requested
                if request.year and r.published.year != request.year:
                    continue
                    
                results.append({
                    "id": r.entry_id.split('/')[-1],
                    "title": r.title,
                    "authors": [a.name for a in r.authors],
                    "year": r.published.year,
                    "abstract": r.summary,
                    "pdf_url": r.pdf_url
                })
            return results
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Arxiv search failed: {e}")

    @router.post("/{discussion_id}/artefacts/arxiv/import", response_model=ArtefactAndDataZoneUpdateResponse)
    async def import_arxiv_selected(
        discussion_id: str,
        request: ArxivImportSelectedRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        import arxiv
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        try:
            for item in request.items:
                search = arxiv.Search(id_list=[item.id])
                paper = next(search.results())
                
                if item.mode == "full":
                    # Download PDF and extract
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
                        paper.download_pdf(filename=tf.name)
                        pdf_path = tf.name
                    
                    pdf_doc = fitz.open(pdf_path)
                    text = "\n".join([page.get_text() for page in pdf_doc])
                    pdf_doc.close()
                    os.unlink(pdf_path)
                    
                    content = f"# {paper.title}\nAuthors: {', '.join([a.name for a in paper.authors])}\nSource: {paper.entry_id}\n\n{text}"
                else:
                    content = f"# {paper.title} (Abstract)\nAuthors: {', '.join([a.name for a in paper.authors])}\nSource: {paper.entry_id}\n\n{paper.summary}"
                
                # Use update_artefact for unified management
                discussion.update_artefact(
                    f"Arxiv_{item.id}.md", 
                    content, 
                    author=current_user.username, 
                    active=request.auto_load
                )
            
            discussion.commit()
            
            artefacts = discussion.list_artefacts()
            for artefact in artefacts:
                if isinstance(artefact.get('created_at'), datetime): artefact['created_at'] = artefact['created_at'].isoformat()
                if isinstance(artefact.get('updated_at'), datetime): artefact['updated_at'] = artefact['updated_at'].isoformat()

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            all_images_info = discussion.get_discussion_images()

            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Arxiv import failed: {e}")

    @router.post("/{discussion_id}/artefacts/github/search", response_model=List[Dict[str, str]])
    async def search_github(
        discussion_id: str,
        request: GithubSearchRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        try:
            query = request.query.strip()
            if query.startswith("http://") or query.startswith("https://"):
                return[{"title": "Direct GitHub URL", "url": query, "snippet": "Direct import."}]
                
            api_url = f"https://api.github.com/search/issues?q={query}&per_page=10"
            r = requests.get(api_url, headers={"Accept": "application/vnd.github.v3+json"})
            r.raise_for_status()
            items = r.json().get('items', [])
            return[{
                "title": i.get('title', 'Unknown'), 
                "url": i.get('html_url', ''), 
                "snippet": f"[{i.get('state', 'unknown').upper()}] {i.get('repository_url', '').split('/')[-1]} | Comments: {i.get('comments', 0)}"
            } for i in items]
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"GitHub search failed: {e}")

    @router.post("/{discussion_id}/artefacts/github", response_model=ArtefactAndDataZoneUpdateResponse)
    async def import_github(
        discussion_id: str,
        request: GithubImportRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        url = request.url.strip()
        content = ""
        title = "Github_Import"
        
        try:
            if "github.com" in url:
                blob_match = re.match(r'https?://github\.com/([^/]+)/([^/]+)/blob/([^/]+)/(.*)', url)
                issue_match = re.match(r'https?://github\.com/([^/]+)/([^/]+)/(issues|pull)/(\d+)', url)
                
                if blob_match:
                    user, repo, branch, filepath = blob_match.groups()
                    raw_url = f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/{filepath}"
                    r = requests.get(raw_url)
                    r.raise_for_status()
                    
                    # Guess language for markdown fence
                    ext = filepath.split('.')[-1] if '.' in filepath else 'txt'
                    content = f"# File: {filepath}\nSource: {url}\n\n```{ext}\n{r.text}\n```"
                    title = f"GH_{filepath.split('/')[-1]}"
                    
                elif issue_match:
                    user, repo, type_, num = issue_match.groups()
                    api_url = f"https://api.github.com/repos/{user}/{repo}/issues/{num}"
                    r = requests.get(api_url, headers={"Accept": "application/vnd.github.v3+json"})
                    r.raise_for_status()
                    data = r.json()
                    
                    title_str = data.get('title', f'{type_.capitalize()} #{num}')
                    body = data.get('body', '')
                    state = data.get('state', 'unknown')
                    content = f"# [{state.upper()}] {title_str}\nSource: {url}\n\n{body}"
                    
                    comments_url = data.get('comments_url')
                    if comments_url:
                        rc = requests.get(comments_url, headers={"Accept": "application/vnd.github.v3+json"})
                        if rc.status_code == 200:
                            comments = rc.json()
                            for c in comments:
                                content += f"\n\n---\n**{c.get('user',{}).get('login','User')}** commented:\n{c.get('body','')}"
                    title = f"GH_{type_}_{num}"
                    
                else:
                    raise HTTPException(status_code=400, detail="Only GitHub file blobs, issues, or pull requests are currently supported.")
                    
            elif "raw.githubusercontent.com" in url:
                r = requests.get(url)
                r.raise_for_status()
                filepath = url.split('/')[-1]
                ext = filepath.split('.')[-1] if '.' in filepath else 'txt'
                content = f"# File: {filepath}\nSource: {url}\n\n```{ext}\n{r.text}\n```"
                title = f"GH_Raw_{filepath}"
            else:
                 raise HTTPException(status_code=400, detail="Not a valid GitHub URL.")

            # Ensure safe filename
            safe_title = re.sub(r'[^A-Za-z0-9_.-]', '_', title) + ".md"
            
            # Use update_artefact for unified management
            discussion.update_artefact(
                safe_title, 
                content, 
                author=current_user.username, 
                active=request.auto_load
            )
            
            discussion.commit()
            
            artefacts = discussion.list_artefacts()
            for artefact in artefacts:
                if isinstance(artefact.get('created_at'), datetime): artefact['created_at'] = artefact['created_at'].isoformat()
                if isinstance(artefact.get('updated_at'), datetime): artefact['updated_at'] = artefact['updated_at'].isoformat()

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            all_images_info = discussion.get_discussion_images()

            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
            
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"GitHub import failed: {e}")

    @router.post("/{discussion_id}/artefacts/stackoverflow/search", response_model=List[Dict[str, str]])
    async def search_stackoverflow(
        discussion_id: str,
        request: StackOverflowSearchRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        try:
            query = request.query.strip()
            if query.startswith("http://") or query.startswith("https://"):
                return[{"title": "Direct StackOverflow URL", "url": query, "snippet": "Direct import."}]

            api_url = f"https://api.stackexchange.com/2.3/search/advanced?order=desc&sort=relevance&q={query}&site=stackoverflow&pagesize=10"
            r = requests.get(api_url)
            r.raise_for_status()
            items = r.json().get('items', [])
            return[{
                "title": i.get('title', 'Unknown'), 
                "url": i.get('link', ''), 
                "snippet": f"Score: {i.get('score', 0)} | Answers: {i.get('answer_count', 0)} | Tags: {', '.join(i.get('tags', []))}"
            } for i in items]
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"StackOverflow search failed: {e}")

    @router.post("/{discussion_id}/artefacts/stackoverflow", response_model=ArtefactAndDataZoneUpdateResponse)
    async def import_stackoverflow(
        discussion_id: str,
        request: StackOverflowImportRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        pm.ensure_packages("markdownify")
        from markdownify import markdownify as md
        
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        url = request.url.strip()
        
        q_match = re.search(r'stackoverflow\.com/questions/(\d+)', url)
        if not q_match:
             raise HTTPException(status_code=400, detail="Invalid StackOverflow question URL. Must contain /questions/ID")
        
        q_id = q_match.group(1)
        
        try:
            # Fetch question
            api_q = f"https://api.stackexchange.com/2.3/questions/{q_id}?site=stackoverflow&filter=withbody"
            r_q = requests.get(api_q)
            r_q.raise_for_status()
            q_data = r_q.json().get('items',[])
            if not q_data:
                raise HTTPException(status_code=404, detail="Question not found on StackOverflow.")
            
            question = q_data[0]
            q_title = question.get('title', 'StackOverflow Question')
            q_body_html = question.get('body', '')
            q_body_md = md(q_body_html).strip()

            content = f"# {q_title}\nSource: {url}\n\n**Question (Score: {question.get('score', 0)}):**\n\n{q_body_md}\n\n---\n"

            # Fetch top 3 answers
            api_a = f"https://api.stackexchange.com/2.3/questions/{q_id}/answers?site=stackoverflow&order=desc&sort=votes&filter=withbody"
            r_a = requests.get(api_a)
            if r_a.status_code == 200:
                a_data = r_a.json().get('items',[])
                for i, ans in enumerate(a_data[:3]):
                    is_accepted = "✅ ACCEPTED " if ans.get('is_accepted') else ""
                    a_body_md = md(ans.get('body', '')).strip()
                    content += f"\n### {is_accepted}Answer {i+1} (Score: {ans.get('score', 0)})\n\n{a_body_md}\n\n---\n"
                
            title = f"SO_{q_id}.md"
            
            # Use update_artefact for unified management
            discussion.update_artefact(
                title, 
                content, 
                author=current_user.username, 
                active=request.auto_load
            )
            
            discussion.commit()
            
            artefacts = discussion.list_artefacts()
            for artefact in artefacts:
                if isinstance(artefact.get('created_at'), datetime): artefact['created_at'] = artefact['created_at'].isoformat()
                if isinstance(artefact.get('updated_at'), datetime): artefact['updated_at'] = artefact['updated_at'].isoformat()

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            all_images_info = discussion.get_discussion_images()

            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }

        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"StackOverflow import failed: {e}")

    @router.post("/{discussion_id}/artefacts/youtube", response_model=ArtefactAndDataZoneUpdateResponse)
    async def import_youtube_transcript(
        discussion_id: str,
        request: YoutubeTranscriptImportRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        if YouTubeTranscriptApi is None:
            raise HTTPException(status_code=501, detail="youtube_transcript_api is not installed on the server.")

        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        # 1. Improved Video ID Extraction
        video_id = None
        patterns = [
            r'(?:v=|\/)([0-9A-Za-z_-]{11}).*', # v=... or /...
            r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})', # youtu.be/ID
            r'(?:embed\/)([0-9A-Za-z_-]{11})' # embed/ID
        ]
        
        for p in patterns:
             match = re.search(p, request.video_url)
             if match:
                 video_id = match.group(1)
                 break
        
        # Fallback: check if the whole string is an ID
        if not video_id:
             if len(request.video_url) == 11 and re.match(r'^[0-9A-Za-z_-]{11}$', request.video_url):
                 video_id = request.video_url
        
        if not video_id:
            raise HTTPException(status_code=400, detail="Could not extract a valid YouTube Video ID from the provided URL.")

        try:
            # 2. Retrieve Transcript List
            try:
                yvt = YouTubeTranscriptApi()
                transcript_list_obj = yvt.list(video_id)
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Failed to retrieve transcript list. Video may not have captions or is restricted. Error: {e}")

            target_transcript = None
            requested_lang = request.language.lower().strip() if request.language else None

            # 3. Smart Selection Logic
            if requested_lang:
                # Try finding exact match
                try:
                    target_transcript = transcript_list_obj.find_transcript([requested_lang])
                except:
                    # Try finding a translation
                    try:
                        # Translate the first available transcript
                        first_available = next(iter(transcript_list_obj))
                        if first_available.is_translatable:
                            target_transcript = first_available.translate(requested_lang)
                    except:
                        pass # Translation failed or not available
            
            # If no specific language requested OR specific lookup failed
            if not target_transcript:
                # Priority: English -> First Available
                try:
                    target_transcript = transcript_list_obj.find_generated_transcript(['en'])
                except:
                    pass
                
                if not target_transcript:
                    try:
                        target_transcript = transcript_list_obj.find_manually_created_transcript(['en'])
                    except:
                        pass
                
                if not target_transcript:
                    try:
                        target_transcript = next(iter(transcript_list_obj))
                    except:
                        pass

            if not target_transcript:
                raise HTTPException(status_code=400, detail="No suitable transcript found.")

            # 4. Fetch
            transcript_data = target_transcript.fetch()

            # 5. Format
            lines = []
            for entry in transcript_data.snippets:
                start = int(entry.start)
                minutes = start // 60
                seconds = start % 60
                text = entry.text
                lines.append(f"[{minutes:02d}:{seconds:02d}] {text}")
            
            lang_label = target_transcript.language if hasattr(target_transcript, 'language') else (requested_lang or 'unknown')
            full_content = f"# YouTube Transcript ({lang_label})\nSource: {request.video_url}\n\n" + "\n".join(lines)
            
            # 6. Save using the internal manager
            artefact_name = f"Youtube_Transcript_{video_id}.md"
            discussion.update_artefact(
                artefact_name,
                full_content,
                author=current_user.username,
                active=request.auto_load
            )
            
            # Activation is already handled by the 'active' parameter logic (default is True in add_artefact)
            
            discussion.commit()
            
            artefacts = discussion.list_artefacts()
            for artefact in artefacts:
                if isinstance(artefact.get('created_at'), datetime): artefact['created_at'] = artefact['created_at'].isoformat()
                if isinstance(artefact.get('updated_at'), datetime): artefact['updated_at'] = artefact['updated_at'].isoformat()

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            all_images_info = discussion.get_discussion_images()

            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }

        except HTTPException as he:
            raise he
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Failed to process YouTube transcript: {str(e)}")

    class RenameArtefactRequest(BaseModel):
        old_title: str
        new_title: str
        new_type: Optional[str] = None

    @router.put("/{discussion_id}/artefacts/rename", status_code=status.HTTP_200_OK)
    async def rename_discussion_artefact(
        discussion_id: str,
        payload: RenameArtefactRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        # 1. Type Auto-Detection from extension
        ext = Path(payload.new_title).suffix.lower()
        auto_type = payload.new_type
        
        if not auto_type:
            CODE_EXTS = {".py", ".js", ".ts", ".html", ".css", ".c", ".cpp", ".h", ".cs", ".java", ".sh", ".sql", ".vhd", ".v", ".rb", ".php", ".go", ".rs", ".swift", ".kt"}
            if ext in CODE_EXTS:
                auto_type = "code"
            elif ext in {".md", ".txt", ".pdf", ".docx", ".xlsx", ".pptx"}:
                auto_type = "document"

        try:
            # We use the library's internal artefacts manager to perform the rename
            # This ensures all versions associated with the old title are migrated.
            discussion.artefacts.rename(
                old_title=payload.old_title, 
                new_title=payload.new_title, 
                new_type=auto_type
            )
            discussion.commit()
            return {"message": f"Artefact renamed to '{payload.new_title}'", "detected_type": auto_type}
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @router.post("/{discussion_id}/artefacts/manual", response_model=ArtefactAndDataZoneUpdateResponse, status_code=status.HTTP_201_CREATED)
    async def create_manual_artefact(
        discussion_id: str,
        payload: ArtefactCreateManual,
        artefact_type: str = Query(None),
        auto_load: bool = Query(True),
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        try:
            # Type resolution logic
            raw_type = (artefact_type or "").lower()
            if "note" in raw_type: final_type = "note"
            elif "skill" in raw_type: final_type = "skill"
            else:
                ext = payload.title.split('.')[-1].lower() if '.' in payload.title else ""
                final_type = "code" if ext in ['py', 'js', 'ts', 'html', 'css', 'sql', 'cpp', 'c', 'sh'] else "document"

            # Use add_artefact for initial manual creation
            artefact_info = discussion.add_artefact(
                payload.title, 
                payload.content, 
                images=payload.images_b64, 
                author=current_user.username,
                active=auto_load,
                artefact_type=final_type
            )
            
            discussion.commit()
            
            # Fetch latest synchronized state
            raw_artefacts = discussion.list_artefacts()
            all_images_info = discussion.get_discussion_images()

            return {
                "artefacts": [_map_artefact_for_ui(art, discussion_id) for art in raw_artefacts],
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Failed to create artefact: {e}")

    @router.put("/{discussion_id}/artefacts/{artefact_title}", response_model=ArtefactInfo)
    async def update_manual_artefact(
        discussion_id: str,
        artefact_title: str,
        payload: ArtefactUpdate,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        
        try:
            # Prepare metadata overrides
            extra_metadata = {}
            if payload.artefact_type:
                extra_metadata['type'] = payload.artefact_type

            # Use positional arguments for required fields (title, new_content)
            # CRITICAL: We remove 'version' from this call to fix the TypeError in the library
            # which occurs when both the router and library internal logic try to set the kwarg.
            artefact_info = discussion.update_artefact(
                artefact_title, 
                payload.new_content, 
                new_images=payload.kept_images_b64 + payload.new_images_b64,
                update_in_place=payload.update_in_place,
                **extra_metadata
            )
            discussion.commit()

            if isinstance(artefact_info.get('created_at'), datetime):
                artefact_info['created_at'] = artefact_info['created_at'].isoformat()
            if isinstance(artefact_info.get('updated_at'), datetime):
                artefact_info['updated_at'] = artefact_info['updated_at'].isoformat()
            
            return artefact_info

        except ValueError as e:
            raise HTTPException(status_code=404, detail=str(e))
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"An unexpected error occurred while updating the artefact: {e}")

    @router.post("/{discussion_id}/artefacts/export-context", response_model=ArtefactInfo)
    async def export_context_as_artefact(
        discussion_id: str,
        request: ExportContextRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        content = discussion.discussion_data_zone
        if not content or not content.strip():
            raise HTTPException(status_code=400, detail="Data zone is empty.")
        try:
            # Use update_artefact for consistent versioned storage of the context snapshot
            artefact_info = discussion.update_artefact(
                request.title, 
                content, 
                author=current_user.username,
                active=True,
                artefact_type="document"
            )
            discussion.commit()
            return _map_artefact_for_ui(artefact_info, discussion_id)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to create artefact from context: {e}")

    @router.get("/{discussion_id}/artefacts/{artefact_title:path}/content")
    async def get_discussion_artefact_content(
        discussion_id: str,
        artefact_title: str,
        version: Optional[int] = Query(None),
        strategy: str = Query("raw"),  # 'raw' or 'formatted'
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        """
        Get the raw content of a specific artefact version.
        Path-style route to avoid conflicts with SPA catch-all.
        """
        from urllib.parse import unquote
        
        discussion, owner_username, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db)
        
        # URL decode the title (handles spaces and special chars)
        decoded_title = unquote(artefact_title)
        
        artefact = discussion.get_artefact(title=decoded_title, version=version)
        if not artefact:
            raise HTTPException(status_code=404, detail=f"Artefact '{decoded_title}' not found")

        # Return raw content as plain text for the workspace editor
        content = artefact.get('content', '')
        
        from fastapi.responses import PlainTextResponse
        return PlainTextResponse(
            content=content,
            media_type="text/plain; charset=utf-8",
            headers={"X-Artefact-Title": decoded_title, "X-Artefact-Version": str(artefact.get('version', 1))}
        )

    # Keep the old endpoint for backward compatibility (returns JSON with metadata)
    @router.get("/{discussion_id}/artefact", response_model=ArtefactInfo)
    async def get_discussion_artefact_info(
        discussion_id: str,
        artefact_title: str = Query(...),
        version: Optional[int] = Query(None),
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        """
        Get artefact metadata (JSON). Use /content endpoint for raw content.
        """
        discussion, owner_username, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db)
        artefact = discussion.get_artefact(title=artefact_title, version=version)

        if not artefact:
            raise HTTPException(status_code=404, detail="Artefact not found")

        if isinstance(artefact.get('created_at'), datetime):
            artefact['created_at'] = artefact['created_at'].isoformat()
        if isinstance(artefact.get('updated_at'), datetime):
            artefact['updated_at'] = artefact['updated_at'].isoformat()
        return artefact

    @router.delete("/{discussion_id}/artefact", status_code=status.HTTP_200_OK)
    async def delete_discussion_artefact(
        discussion_id: str,
        artefact_title: str = Query(...),
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        from urllib.parse import unquote
        title = unquote(artefact_title) # Ensure URL encoded titles match correctly
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        try:
            # Check if it exists first
            if not discussion.get_artefact(title=title):
                raise HTTPException(status_code=404, detail="Artefact not found.")
                
            discussion.remove_artefact(title=title)
            discussion.commit()
            return {"message": f"Artefact '{title}' deleted."}
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=str(e))

    @router.post("/{discussion_id}/artefacts/revert", status_code=status.HTTP_200_OK)
    async def revert_discussion_artefact(
        discussion_id: str,
        title: str = Body(..., embed=True),
        version: int = Body(..., embed=True),
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        try:
            discussion.artefacts.revert(title, target_version=version)
            discussion.commit()
            return {"message": f"Reverted to version {version}"}
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Revert failed: {e}")

    @router.post("/{discussion_id}/artefacts/load-all-to-context", response_model=ArtefactAndDataZoneUpdateResponse)
    async def load_all_artefacts_to_data_zone(
        discussion_id: str,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        try:
            # The library manages the context layer. We just activate everything.
            all_artefacts_infos = discussion.list_artefacts()
            
            for art_info in all_artefacts_infos:
                discussion.artefacts.activate(art_info['title'], version=art_info['version'])

            discussion.commit()
            
            raw_artefacts = discussion.list_artefacts()
            artefacts = [_map_artefact_for_ui(art) for art in raw_artefacts]

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))
            
            all_images_info = discussion.get_discussion_images()
            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except ValueError as e:
            raise HTTPException(status_code=404, detail=str(e))
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to load all artefacts: {e}")        

    @router.post("/{discussion_id}/artefacts/load-to-context", response_model=ArtefactAndDataZoneUpdateResponse)
    async def load_artefact_to_data_zone(
        discussion_id: str,
        request: LoadArtefactRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        try:
            # We only update the library metadata. The library handles 
            # context injection natively during chat() without modifying the data_zone string.
            discussion.artefacts.activate(request.title, version=request.version)
            discussion.commit()
            
            raw_artefacts = discussion.list_artefacts()
            artefacts = [_map_artefact_for_ui(art) for art in raw_artefacts]
            
            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))

            all_images_info = discussion.get_discussion_images()
            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except ValueError as e:
            raise HTTPException(status_code=404, detail=str(e))
        except Exception as e:
            trace_exception(e)
            raise HTTPException(status_code=500, detail=f"Failed to load artefact: {e}")

    @router.post("/{discussion_id}/artefacts/unload-from-context", response_model=ArtefactAndDataZoneUpdateResponse)
    async def unload_artefact_from_data_zone(
        discussion_id: str,
        request: UnloadArtefactRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        discussion, _, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        try:
            # Deactivate in the library metadata only.
            discussion.artefacts.deactivate(request.title, version=request.version)
            discussion.commit()
            
            raw_artefacts = discussion.list_artefacts()
            artefacts = [_map_artefact_for_ui(art) for art in raw_artefacts]

            lc = get_user_lollms_client(current_user.username)
            token_count = len(lc.tokenize(discussion.discussion_data_zone))

            all_images_info = discussion.get_discussion_images()
            # Pure Artefact Response: We no longer return or touch the discussion_data_zone text
            return {
                "artefacts": artefacts,
                "discussion_images": [img['data'] for img in all_images_info],
                "active_discussion_images": [img['active'] for img in all_images_info]
            }
        except ValueError as e:
            raise HTTPException(status_code=404, detail=str(e))
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to unload artefact: {e}")

    @router.post("/{discussion_id}/artefacts/import_url", response_model=TaskInfo, status_code=status.HTTP_202_ACCEPTED)
    async def import_artefact_from_url(
        discussion_id: str,
        request: UrlImportRequest,
        current_user: UserAuthDetails = Depends(get_current_active_user),
        db: Session = Depends(get_db)
    ):
        _, owner_username, _, _ = await get_discussion_and_owner_for_request(discussion_id, current_user, db, 'interact')
        task = task_manager.submit_task(
            name=f"Importing artefact from URL: {request.url}",
            target=_import_artefact_from_url_task,
            args=(owner_username, discussion_id, request.url, request.depth, request.process_with_ai),
            description=f"Scraping content (depth {request.depth}) and saving as artefact.",
            owner_username=current_user.username
        )
        return task
