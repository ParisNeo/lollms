# backend/routers/files.py
import base64
import io
import re
import shutil
import uuid
import os
import tempfile
import requests
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional

from fastapi import APIRouter, Depends, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse, Response
from werkzeug.utils import secure_filename
import markdown2 
from bs4 import BeautifulSoup 
from docx import Document as DocxDocument 
from docx.shared import Pt, Inches
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation 
from pptx.util import Inches as PptxInches, Pt as PptxPt
from markdown_pdf import MarkdownPdf, Section 
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from latex2mathml.converter import convert as latex2mathml


# Try to import optional document parsing libraries
try:
    from docx2python import docx2python
except ImportError:
    docx2python = None
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    MSO_SHAPE_TYPE = None
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
try:
    import fitz 
except ImportError:
    fitz = None
try:
    from mdtopptx import parse_markdown as md_to_pptx_parse, create_ppt
except ImportError:
    md_to_pptx_parse = None
    create_ppt = None


from backend.session import (
    get_current_active_user, get_user_temp_uploads_path
)
from backend.models import UserAuthDetails
from backend.models.files import ContentExportRequest, ExtractionAndEmbeddingResponse
from backend.discussion import get_user_discussion
from ascii_colors import trace_exception
from backend.config import TEMP_UPLOADS_DIR_NAME
from backend.settings import settings

files_router = APIRouter(prefix="/api/files", tags=["Files"])
upload_router = APIRouter(prefix="/api/upload", tags=["Files"])
assets_router = APIRouter(prefix="/assets", tags=["Files"])

def _process_msg_attachment(att_bytes: bytes, att_name: str, images: List[str], extract_images: bool = True) -> Optional[str]:
    """Helper to process a single attachment from an MSG file."""
    att_ext = Path(att_name).suffix.lower()
    
    if att_ext in [".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"]:
        if extract_images:
            images.append(base64.b64encode(att_bytes).decode("utf-8"))
        return None 
    
    try:
        text_guess = att_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text_guess = att_bytes.decode("latin-1", errors="replace")
        except Exception:
            text_guess = ""
    
    CODE_EXTENSIONS = {
        ".txt": "text", ".md": "markdown", ".py": "python", ".js": "javascript", ".ts": "typescript", ".html": "html", ".css": "css",
        ".c": "c", ".cpp": "cpp", ".h": "cpp", ".hpp": "cpp", ".cs": "csharp", ".java": "java",
        ".json": "json", ".xml": "xml", ".sh": "bash", ".vhd": "vhdl", ".v": "verilog",
        ".rb": "ruby", ".php": "php", ".go": "go", ".rs": "rust", ".swift": "swift", ".kt": "kotlin",
        ".yaml": "yaml", ".yml": "yaml", ".sql": "sql", ".log": "text", ".csv": "csv"
    }
    
    if att_ext in CODE_EXTENSIONS:
        lang = CODE_EXTENSIONS[att_ext]
        return f"### Attachment: {att_name}\n\n````{lang}\n{text_guess}\n````"
    
    return f"- Attachment: {att_name} ({len(att_bytes)} bytes - content ignored)"


def extract_text_from_file_bytes(file_bytes: bytes, filename: str, extract_images: bool = True) -> Tuple[str, List[str]]:
    """
    Extracts text and embedded/generated images (as base64) from file bytes.
    Returns: (extracted_text, list_of_base64_images)
    """
    extension = Path(filename).suffix.lower()
    
    extracted_text = ""
    images: List[str] = []
    
    # --- Document Type Handling ---
    
    if extension == ".pdf" and fitz:
        try:
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf_doc:
                text_parts = []
                image_count = 0
                for page in pdf_doc:
                    text_parts.append(page.get_text())
                    if extract_images:
                        img_list = page.get_images(full=True)
                        image_count += len(img_list)
                        for img_info in img_list:
                            xref = img_info[0]
                            base_image = pdf_doc.extract_image(xref)
                            images.append(base64.b64encode(base_image["image"]).decode('utf-8'))
                extracted_text = "\n".join(text_parts).strip()
                
                if not extracted_text and image_count > 0:
                    raise ValueError("This appears to be a scanned PDF with no text layer. LoLLMs cannot process it directly. Please use an OCR (Optical Character Recognition) tool to convert it to a text-based PDF or extract the text manually before uploading.")

        except ValueError as e:
             raise e
        except Exception as e:
            extracted_text = f"[Error processing PDF file: {e}. Is PyMuPDF (fitz) installed?]"
            
    elif extension == ".docx" and docx2python:
        try:
            with io.BytesIO(file_bytes) as docx_io:
                result = docx2python(docx_io)
                extracted_text = result.text
                if result.images and extract_images:
                    for image_bytes in result.images.values():
                        images.append(base64.b64encode(image_bytes).decode("utf-8"))
        except Exception as e:
            extracted_text = f"[Error processing DOCX file: {e}. Is docx2python installed?]"
            
    elif (extension == ".xlsx" or "spreadsheetml" in filename.lower()) and pd:
        try:
            xls = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            md_parts = []
            for sheet_name, df in xls.items():
                md_parts.append(f"### {sheet_name}\n\n{df.to_markdown(index=False)}")
            extracted_text = "\n\n".join(md_parts)
        except Exception as e:
            extracted_text = f"[Error processing XLSX file: {e}. Is pandas installed?]"

    elif extension == ".pptx" and MSO_SHAPE_TYPE:
        try:
            with io.BytesIO(file_bytes) as pptx_io:
                prs = Presentation(pptx_io)
                slide_texts: List[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    slide_parts: List[str] = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt = (shape.text or "").strip()
                            if txt: slide_parts.append(txt)
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and extract_images:
                            images.append(base64.b64encode(shape.image.blob).decode("utf-8"))
                    if slide_parts: slide_texts.append(f"--- Slide {idx} ---\n" + "\n".join(slide_parts))
                extracted_text = "\n\n".join(slide_texts)
        except Exception as e:
            extracted_text = f"[Error processing PPTX file: {e}]"
            
    elif extension == ".msg" and extract_msg:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tf:
                tf.write(file_bytes)
                temp_path = tf.name
            try:
                msg = extract_msg.Message(temp_path)
                header_lines = []
                if getattr(msg, "subject", ""): header_lines.append(f"# {getattr(msg, 'subject', '')}")
                meta = []
                if getattr(msg, "sender", None) or getattr(msg, "from_", None): meta.append(f"From: {getattr(msg, 'sender', None) or getattr(msg, 'from_', None)}")
                if getattr(msg, "to", ""): meta.append(f"To: {getattr(msg, 'to', '')}")
                if getattr(msg, "date", ""): meta.append(f"Date: {getattr(msg, 'date', '')}")
                if meta: header_lines.append("\n".join(meta))
                header = "\n\n".join(header_lines)

                msg_body = (getattr(msg, "body", "") or "").strip()
                if not msg_body and getattr(msg, "htmlBody", None):
                    msg_body = BeautifulSoup(getattr(msg, "htmlBody"), "html.parser").get_text()

                attachment_text_parts: List[str] = [header, msg_body]
                for att in msg.attachments:
                    text_part = _process_msg_attachment(att.data or b"", att.longFilename or att.shortFilename or "attachment", images, extract_images=extract_images)
                    if text_part: attachment_text_parts.append(text_part)

                extracted_text = "\n\n".join([p for p in attachment_text_parts if p.strip()])

            finally:
                os.unlink(temp_path)
        except Exception as e:
            extracted_text = f"[Error processing MSG file: {e}. Is extract_msg installed?]"

    # 6. Plain Text / Code
    else:
        try:
            extracted_text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            extracted_text = file_bytes.decode('latin-1', errors='replace')

    # Add markdown code fence for recognized code files
    CODE_EXTENSIONS = {
        ".py", ".js", ".ts", ".html", ".css", ".c", ".cpp", ".h", ".hpp", ".cs", ".java",
        ".json", ".xml", ".sh", ".vhd", ".v", ".rb", ".php", ".go", ".rs", ".swift", ".kt",
        ".yaml", ".yml", ".sql", ".log", ".csv", ".txt", ".md"
    }
    if extension in CODE_EXTENSIONS:
        lang = extension.strip('.').replace('c++', 'cpp').replace('c#', 'csharp')
        if not extracted_text.startswith('```'):
            extracted_text = f"````{lang}\n{extracted_text}\n````"
            
    return extracted_text, images

@files_router.post("/extract-text")
async def extract_text_from_file(
    file: UploadFile = File(...)
):
    """
    Extracts text content from a single uploaded file.
    Supports various formats like PDF, DOCX, TXT, etc.
    """
    try:
        content_bytes = await file.read()
        text_content, _ = extract_text_from_file_bytes(content_bytes, file.filename, extract_images=False)
        return {"text_content": text_content}
    except Exception as e:
        trace_exception(e)
        raise HTTPException(status_code=500, detail=f"Failed to extract text from file: {str(e)}")

@files_router.post("/export-markdown")
async def export_as_markdown(
    content: str = Form(...),
    filename: str = Form("export.md")
):
    """
    Accepts text content and returns it as a downloadable Markdown file.
    """
    safe_filename = secure_filename(filename)
    if not safe_filename.endswith('.md'):
        safe_filename += '.md'
    
    headers = {
        'Content-Disposition': f'attachment; filename="{safe_filename}"'
    }
    return Response(content=content, media_type='text/markdown', headers=headers)

# Assume 'settings' and 'ContentExportRequest' and trace_exception exist in your module

def md2_to_html(md_text: str) -> str:
    # Configure markdown2 extras to cover requested features
    extras = [
        "fenced-code-blocks",
        "tables",
        "footnotes",
        "smarty-pants",
        "strike",
        "spoiler",
        "code-friendly",
        "cuddled-lists",
        "header-ids",
        "toc"
    ]
    return markdown2.markdown(md_text, extras=extras)

def _download_image_to_temp(src: str) -> str:
    if src.startswith("data:"):
        head, b64data = src.split(",", 1)
        suffix = ".png"
        if ";base64" in head and "image/" in head:
            mime = head.split(":")[1].split(";")[0]
            ext = "." + mime.split("/")[1]
            suffix = ext
        data = base64.b64decode(b64data)
        tf = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tf.write(data)
        tf.flush(); tf.close()
        return tf.name
    r = requests.get(src, timeout=10)
    r.raise_for_status()
    ext = os.path.splitext(src)[1] or ".png"
    tf = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    tf.write(r.content)
    tf.flush(); tf.close()
    return tf.name

def _insert_math(p, latex):
    mathml = latex2mathml(latex)
    # Replace MathML namespace for Office Math
    omath_xml = (
        f'<m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">'
        f'<m:oMath>'
        f'{mathml}'
        f'</m:oMath>'
        f'</m:oMathPara>'
    )
    # Attach OMML to paragraph
    omath_element = parse_xml(omath_xml)
    p._p.append(omath_element)
    
def html_to_docx_bytes(html: str) -> bytes:
    soup = BeautifulSoup(html, "html.parser")
    doc = DocxDocument()

    def add_inline_runs(p, node):
        for child in node.children:
            if isinstance(child, str):
                # Detect and handle inline math: \( ... \)
                parts = re.split(r'(\\\(.+?\\\)|\\\[.+?\\\])', child)
                for part in parts:
                    if part.startswith(r'\(') and part.endswith(r'\)'):
                        latex = part[2:-2].strip()
                        _insert_math(p, latex)
                    elif part.startswith(r'\[') and part.endswith(r'\]'):
                        latex = part[2:-2].strip()
                        _insert_math(p, latex)
                    else:
                        p.add_run(part)
                continue
            tag = getattr(child, "name", "").lower()
            if tag in ("strong", "b"):
                run = p.add_run(child.get_text()); run.bold = True
            elif tag in ("em", "i"):
                run = p.add_run(child.get_text()); run.italic = True
            elif tag == "code":
                run = p.add_run(child.get_text()); run.font.name = "Consolas"; run.font.size = Pt(10)
            elif tag == "a":
                text = child.get_text(); href = child.get("href") or ""
                run = p.add_run(text + f" ({href})"); run.underline = True
            else:
                add_inline_runs(p, child)

    def handle_block(el):
        if not getattr(el, "name", None):
            return
        name = el.name.lower()

        if name in [f"h{i}" for i in range(1, 7)]:
            lvl = int(name[1]); doc.add_heading(el.get_text(), level=min(9, lvl)); return

        if name == "p":
            p = doc.add_paragraph(); add_inline_runs(p, el); return

        if name in ("ul", "ol"):
            style = "List Bullet" if name == "ul" else "List Number"
            for li in el.find_all("li", recursive=False):
                p = doc.add_paragraph(style=style); add_inline_runs(p, li)
            return

        if name == "blockquote":
            p = doc.add_paragraph(style="Intense Quote"); add_inline_runs(p, el); return

        if name == "pre":
            code_text = el.get_text()
            p = doc.add_paragraph(); run = p.add_run(code_text)
            run.font.name = "Consolas"; run.font.size = Pt(10); return

        if name == "table":
            rows = el.find_all("tr")
            if not rows: return
            cols = rows[0].find_all(["th", "td"])
            table = doc.add_table(rows=len(rows), cols=len(cols))
            for i, row in enumerate(rows):
                cells = row.find_all(["th", "td"])
                for j, cell in enumerate(cells):
                    table.cell(i, j).text = cell.get_text()
            return

        if name == "img":
            src = el.get("src"); alt = el.get("alt") or "[image]"
            tmp = None
            try:
                tmp = _download_image_to_temp(src)
                doc.add_picture(tmp, width=Inches(5.5))
            except Exception:
                doc.add_paragraph(alt)
            finally:
                try:
                    if tmp and os.path.exists(tmp): os.unlink(tmp)
                except Exception:
                    pass
            return

        # Recurse into generic containers
        for child in el.children:
            handle_block(child)

    body = soup.body or soup
    for child in body.children:
        handle_block(child)

    bio = io.BytesIO(); doc.save(bio); return bio.getvalue()

def md_to_pdf_bytes(md_text: str, toc_level: int = 3) -> bytes:
    pdf = MarkdownPdf(toc_level=toc_level)
    pdf.add_section(Section(md_text))
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
        out_path = tf.name
    pdf.save(out_path)
    with open(out_path, "rb") as f:
        data = f.read()
    try: os.unlink(out_path)
    except Exception: pass
    return data

def md_to_pptx_bytes(md_text: str) -> bytes:
    """
    Robust Markdown to PowerPoint converter with improved text chunking and image support.
    
    1. Extracts images embedded in markdown (e.g. ![alt](data:image...)).
    2. Removes image tags from text to process separately.
    3. Splits remaining text into logical chunks (headers, paragraphs) to fit slides.
    4. Creates separate slides for images.
    """
    prs = Presentation()
    
    # 1. Extract Images
    # Regex to find markdown images: ![alt](url)
    image_pattern = re.compile(r'!\[.*?\]\((.*?)\)')
    found_images = image_pattern.findall(md_text)
    
    # 2. Clean Text
    # Remove image tags to leave only text for chunking
    clean_text = image_pattern.sub('', md_text).strip()
    
    # --- Helper: Add Title Slide ---
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title if title else "Presentation"
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    # --- Helper: Add Content Slide ---
    def add_content_slide(title, body_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        # Basic markdown cleanup for body
        # Convert bullets
        clean_body = body_text.replace('* ', '• ').replace('- ', '• ')
        
        tf = slide.placeholders[1].text_frame
        tf.text = clean_body

    # --- Helper: Add Image Slide ---
    def add_image_slide(image_src):
        # Blank layout for images to maximize space
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        tmp_path = None
        try:
            tmp_path = _download_image_to_temp(image_src)
            
            # Add picture
            # Center and fit logic
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # We add it first to get dimensions, then adjust
            pic = slide.shapes.add_picture(tmp_path, 0, 0)
            
            # Scale to fit within slide
            # Calculate ratios
            image_ratio = pic.width / pic.height
            slide_ratio = slide_width / slide_height
            
            if image_ratio > slide_ratio:
                # Image is wider relative to slide -> fit width
                pic.width = slide_width
                pic.height = int(slide_width / image_ratio)
                pic.left = 0
                pic.top = int((slide_height - pic.height) / 2)
            else:
                # Image is taller relative to slide -> fit height
                pic.height = slide_height
                pic.width = int(slide_height * image_ratio)
                pic.top = 0
                pic.left = int((slide_width - pic.width) / 2)

        except Exception as e:
            print(f"Failed to add image to slide: {e}")
            txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(1), PptxInches(8), PptxInches(5))
            tf = txBox.text_frame
            tf.text = "[Image could not be loaded]"
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try: os.unlink(tmp_path)
                except: pass

    # --- Logic: Create Slides ---
    
    # 1. Image Slides (If any exist, we put them first or interspersed? 
    # Current request: "use those images as the slides". We'll put them first for visibility, 
    # or if text refers to them, interspersing is hard without complex parsing.
    # Let's add them at the end of the presentation to not obscure the intro text, OR
    # if the text is empty, just images.
    # A safe bet for "slide deck" generation is Intro -> Text Content -> Images (Appendix/Gallery).
    
    # 3. Intelligent Text Chunking
    chunks = []
    
    # Split by explicit horizontal rules first
    segments = re.split(r'\n---\n', clean_text)
    
    for segment in segments:
        if not segment.strip(): continue
        
        # Split by Headers (#)
        header_parts = re.split(r'\n(#+ )', '\n' + segment)
        # header_parts will look like ['', '# ', 'Title', '## ', 'Subtitle...', ...]
        
        current_title = "Slide"
        current_body = ""
        
        # Iterate and reconstruct
        i = 0
        while i < len(header_parts):
            part = header_parts[i]
            
            if part.strip().startswith('#'):
                # It's a header marker, next part is the content line
                if i + 1 < len(header_parts):
                    # Extract title line
                    full_line = header_parts[i+1].split('\n', 1)
                    current_title = full_line[0].strip()
                    
                    # Any content after the title on the same chunk
                    rest_of_section = full_line[1] if len(full_line) > 1 else ""
                    
                    if current_body.strip():
                        chunks.append({'title': "Content", 'body': current_body.strip()})
                        current_body = ""
                    
                    current_body = rest_of_section
                    i += 2
                else:
                    i += 1
            else:
                current_body += part
                i += 1
        
        if current_body.strip():
            # Check length of body
            # PPTX slide fits roughly 600-800 chars comfortably depending on font size
            MAX_CHARS = 700
            if len(current_body) > MAX_CHARS:
                # Split by paragraphs
                paras = current_body.split('\n\n')
                temp_chunk = ""
                for para in paras:
                    if len(temp_chunk) + len(para) < MAX_CHARS:
                        temp_chunk += para + "\n\n"
                    else:
                        chunks.append({'title': current_title, 'body': temp_chunk.strip()})
                        temp_chunk = para + "\n\n"
                if temp_chunk.strip():
                     chunks.append({'title': current_title, 'body': temp_chunk.strip()})
            else:
                chunks.append({'title': current_title, 'body': current_body.strip()})

    # Add Text Slides
    for chunk in chunks:
        add_content_slide(chunk['title'], chunk['body'])

    # Add Image Slides
    for img_src in found_images:
        add_image_slide(img_src)
        
    # If no content at all, add a placeholder
    if len(prs.slides) == 0:
        add_title_slide("Empty Presentation")

    # Output
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

def html_wrapper(html_body: str, title: str = "Export") -> bytes:
    return f"<html><head><meta charset='utf-8'><title>{title}</title></head><body>{html_body}</body></html>".encode("utf-8")

@files_router.post("/export-content")
async def export_content(payload: ContentExportRequest):
    export_format = payload.format.lower()
    setting_key_format = 'markdown' if export_format == 'md' else export_format
    setting_key = f"export_to_{setting_key_format}_enabled"

    if not settings.get(setting_key, False):
        raise HTTPException(status_code=403, detail=f"Export to '{export_format}' is disabled by the administrator.")

    content = payload.content
    filename = f"export.{export_format}"
    media_type = "application/octet-stream"
    file_content = b''

    try:
        if export_format == 'txt':
            media_type = "text/plain"
            file_content = content.encode('utf-8')

        elif export_format in ['md', 'markdown']:
            media_type = "text/markdown"
            file_content = content.encode('utf-8')

        elif export_format == 'html':
            html_content = md2_to_html(content)  # markdown2
            media_type = "text/html"
            file_content = html_wrapper(html_content, title="Export")  # consistent HTML shell

        elif export_format == 'pdf':
            media_type = "application/pdf"
            file_content = md_to_pdf_bytes(content)  # preserves headings/lists/code/images/tables

        elif export_format == 'docx':
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            html_content = md2_to_html(content)  # markdown2
            file_content = html_to_docx_bytes(html_content)  # map HTML → Word

        elif export_format == 'pptx':
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            file_content = md_to_pptx_bytes(content)  # Markdown slides via custom logic

        elif export_format == 'xlsx':
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            wb = Workbook()
            table_regex = re.compile(r'(\|.*\|(?:\r?\n|\r)?\|[-| :]*\|(?:\r?\n|\r)?(?:\|.*\|(?:\r?\n|\r)?)+)')
            tables = table_regex.findall(content)

            if tables:
                for i, table_md in enumerate(tables):
                    ws = wb.create_sheet(title=f"Table {i+1}")
                    lines = [line.strip() for line in table_md.strip().split('\n')]
                    data_rows = [line for line in lines if not re.match(r'^\|[-| :]*\|$', line)]
                    for row_idx, line in enumerate(data_rows):
                        cells = [cell.strip() for cell in line.strip('|').split('|')]
                        for col_idx, cell_text in enumerate(cells):
                            ws.cell(row=row_idx + 1, column=col_idx + 1, value=cell_text)
                    for col in ws.columns:
                        max_length = 0
                        column = get_column_letter(col[0].column)
                        for cell in col:
                            if cell.value and len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        ws.column_dimensions[column].width = (max_length + 2)
                if wb.sheetnames[0] == 'Sheet':
                    wb.remove(wb['Sheet'])
            else:
                ws = wb.active
                ws.title = "Content"
                ws['A1'] = content
                ws.column_dimensions['A'].width = 100
                ws['A1'].alignment = ws['A1'].alignment.copy(wrapText=True)

            bio = io.BytesIO(); wb.save(bio); file_content = bio.getvalue()

        elif export_format == 'epub':
            media_type = "application/epub+zip"
            html_content = md2_to_html(content)
            file_content = html_wrapper(html_content, title="Export")  # placeholder xhtml; replace with proper EPUB packaging

        elif export_format == 'odt':
            media_type = "application/vnd.oasis.opendocument.text"
            html_content = md2_to_html(content)
            file_content = html_wrapper(html_content, title="Export")  # placeholder

        elif export_format == 'rtf':
            media_type = "application/rtf"
            text = BeautifulSoup(md2_to_html(content), "html.parser").get_text()
            rtf = r"{\rtf1\ansi " + text.replace("\\", r"\\").replace("{", r"\{").replace("}", r"\}").replace("\n", r"\par ") + "}"
            file_content = rtf.encode("utf-8", errors="ignore")

        elif export_format in ['tex', 'latex']:
            media_type = "application/x-tex"
            text = BeautifulSoup(md2_to_html(content), "html.parser").get_text()
            def esc(t: str) -> str:
                rep = {"\\": r"\textbackslash{}", "&": r"\&", "%": r"\%", "$": r"\$", "#": r"\#", "_": r"\_", "{": r"\{", "}": r"\}", "~": r"\textasciitilde{}", "^": r"\textasciicircum{}"}
                for k, v in rep.items(): t = t.replace(k, v)
                return t
            tex = "\\documentclass{article}\n\\usepackage[utf8]{inputenc}\n\\begin{document}\n" + esc(text) + "\n\\end{document}\n"
            file_content = tex.encode("utf-8")

        else:
            raise HTTPException(status_code=400, detail="Unsupported export format")

    except ImportError as e:
        raise HTTPException(status_code=501, detail=f"A required library for '{export_format}' export is missing: {e.name}")
    except Exception as e:
        trace_exception(e)
        raise HTTPException(status_code=500, detail=f"Failed to generate export file: {e}")

    headers = {'Content-Disposition': f'attachment; filename=\"{filename}\"'}
    return Response(content=file_content, media_type=media_type, headers=headers)

@files_router.post("/extract_and_embed/{discussion_id}", response_model=ExtractionAndEmbeddingResponse)
async def extract_and_embed_files(
    discussion_id: str,
    files: List[UploadFile] = File(...),
    extract_images: bool = Form(True),
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    discussion = get_user_discussion(current_user.username, discussion_id)
    if not discussion:
        raise HTTPException(status_code=404, detail="Discussion not found.")

    all_extracted_text = ""
    initial_image_count = len(discussion.images)
    new_image_count = 0

    temp_img_root = get_user_temp_uploads_path(current_user.username)

    for file in files:
        filename = secure_filename(file.filename)
        content = await file.read()
        
        file_text = f"\n\n--- Document: {filename} ---\n"
        
        try:
            extracted_text, images = extract_text_from_file_bytes(content, filename, extract_images=extract_images)
            
            # --- Integrate the extracted data into the discussion ---
            file_text += extracted_text
            
            if images:
                for b64_image in images:
                    discussion.add_discussion_image(b64_image)
                    new_image_count += 1
                    # Append a reference tag to the text if an image was extracted
                    file_text += f"\n![Image from {filename} | Image {initial_image_count + new_image_count}]\n"
            # --- End Integration ---

            file_text += f"\n--- End Document: {filename} ---\n"
            all_extracted_text += file_text

        except Exception as e:
            trace_exception(e)            
            all_extracted_text += f"\n\n--- Error processing {filename}: {str(e)} ---\n\n"

    discussion.discussion_data_zone = all_extracted_text + (discussion.discussion_data_zone or "")
    discussion.commit()
    all_images_info = discussion.get_discussion_images()

    return ExtractionAndEmbeddingResponse(
        text_content=discussion.discussion_data_zone,
        discussion_images=[img_info['data'] for img_info in all_images_info],
        active_discussion_images=[img_info['active'] for img_info in all_images_info]
    )

@upload_router.post("/chat_image", response_model=List[Dict[str, str]])
async def upload_chat_image(
    files: List[UploadFile] = File(...),
    current_user: UserAuthDetails = Depends(get_current_active_user)
):
    """Handles image uploads specifically for attaching to a chat message before sending."""
    temp_path = get_user_temp_uploads_path(current_user.username)
    uploaded_files = []
    for file in files:
        if not file.content_type.startswith("image/"):
            raise HTTPException(status_code=400, detail="Only image files are allowed.")
        
        s_filename = secure_filename(file.filename)
        # Use a more unique name to avoid collisions
        unique_filename = f"{uuid.uuid4().hex[:8]}_{s_filename}"
        file_path = temp_path / unique_filename
        
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        uploaded_files.append({"filename": s_filename, "server_path": unique_filename})
        
    return uploaded_files
